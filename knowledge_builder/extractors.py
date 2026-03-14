from __future__ import annotations

import csv
import json
import re
from contextlib import contextmanager
from pathlib import Path
from typing import Callable
from xml.etree import ElementTree

from .models import Config, ExtractedContent, PageText
from .utils import detect_date_from_text, word_count


ExtractorFn = Callable[[Path, Config | None, str | None], ExtractedContent]
TEXT_ENCODINGS = ("utf-8", "utf-8-sig", "cp1252", "latin-1")
MAX_PREVIEW_CHARS = 900
LOW_TEXT_WORDS = 25


def get_supported_doc_type(path: Path) -> str | None:
    suffix = path.suffix.lower().lstrip(".")
    if suffix in {"pdf", "docx", "txt", "md", "html", "htm", "xlsx", "pptx", "csv", "tsv", "xml", "json", "png", "jpg", "jpeg"}:
        return "html" if suffix == "htm" else suffix
    if suffix in {"yaml", "yml", "ini", "cfg", "conf", "log", "rst", "toml", "rtf"}:
        return "txt"
    return None


def extract(path: Path, doc_type: str, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    extractors: dict[str, ExtractorFn] = {
        "txt": _extract_text,
        "md": _extract_text,
        "csv": _extract_csv,
        "tsv": _extract_csv,
        "html": _extract_html,
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "xlsx": _extract_xlsx,
        "pptx": _extract_pptx,
        "xml": _extract_xml,
        "json": _extract_json,
        "png": _extract_image,
        "jpg": _extract_image,
        "jpeg": _extract_image,
    }
    fn = extractors.get(doc_type)
    if not fn:
        return _finalize_content(
            path,
            ExtractedContent(
                text="",
                title=path.stem,
                extraction_method="unsupported",
                extraction_status="unsupported",
                failure_reason=f"Unsupported document type: {doc_type}",
                warnings=[f"Unsupported document type: {doc_type}"],
                quality_score=0.0,
            ),
        )

    try:
        content = fn(path, config, strategy)
    except Exception as exc:
        content = ExtractedContent(
            text="",
            title=path.stem,
            extraction_method=f"{doc_type}-failed",
            extraction_status="failed",
            failure_reason=str(exc),
            warnings=[f"Extractor raised an exception: {exc}"],
            quality_score=0.0,
        )
    return _finalize_content(path, content)


def _finalize_content(path: Path, content: ExtractedContent) -> ExtractedContent:
    if not content.title:
        content.title = path.stem
    if not content.doc_date:
        content.doc_date = detect_date_from_text(content.text[:5000]) or None
    if not content.pages and content.text:
        content.pages = [PageText(page_number=1, text=content.text, ocr_used=content.ocr_used)]
    content.preview_excerpt = _preview_excerpt(content)
    if not content.fallback_chain:
        content.fallback_chain = [content.extraction_method]
    content.quality_score = max(0.0, min(1.0, _estimate_quality_score(content)))
    if content.extraction_status == "success" and word_count(content.text) < LOW_TEXT_WORDS and not content.failure_reason:
        content.extraction_status = "partial"
        content.warnings.append(f"Low text yield after extraction ({word_count(content.text)} words).")
    return content


def _estimate_quality_score(content: ExtractedContent) -> float:
    status_scores = {
        "success": 1.0,
        "partial": 0.6,
        "metadata_only": 0.25,
        "unsupported": 0.1,
        "failed": 0.0,
    }
    base = status_scores.get(content.extraction_status, 0.5)
    if content.ocr_used:
        base -= 0.1
    if content.warnings:
        base -= min(0.25, len(content.warnings) * 0.05)
    if word_count(content.text) < LOW_TEXT_WORDS:
        base -= 0.15
    if content.failure_reason:
        base -= 0.2
    return base


def _preview_excerpt(content: ExtractedContent) -> str:
    if content.text.strip():
        return content.text.strip().replace("\r", "")[:MAX_PREVIEW_CHARS]
    if content.failure_reason:
        return content.failure_reason[:MAX_PREVIEW_CHARS]
    if content.warnings:
        return " | ".join(content.warnings)[:MAX_PREVIEW_CHARS]
    return ""


def _make_metadata_only(path: Path, extraction_method: str, warning: str, *, failure_reason: str | None = None) -> ExtractedContent:
    return ExtractedContent(
        text="",
        title=path.stem,
        extraction_method=extraction_method,
        extraction_status="metadata_only",
        warnings=[warning],
        failure_reason=failure_reason,
        quality_score=0.25,
    )


def _read_text_best_effort(path: Path) -> tuple[str, list[str]]:
    warnings: list[str] = []
    for index, encoding in enumerate(TEXT_ENCODINGS):
        try:
            text = path.read_text(encoding=encoding)
            if index > 0:
                warnings.append(f"Decoded with fallback encoding {encoding}.")
            return text, warnings
        except UnicodeDecodeError:
            continue
    text = path.read_text(encoding="utf-8", errors="ignore")
    warnings.append("Decoded with utf-8 ignore fallback due to encoding issues.")
    return text, warnings


def _extract_text(path: Path, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    text, warnings = _read_text_best_effort(path)
    return ExtractedContent(
        text=text,
        title=path.stem,
        extraction_method="text",
        pages=[PageText(page_number=1, text=text)],
        warnings=warnings,
    )


def _extract_csv(path: Path, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    rows: list[str] = []
    warnings: list[str] = []
    try:
        with path.open("r", encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.reader(f, delimiter="\t" if path.suffix.lower() == ".tsv" else ",")
            for row in reader:
                rows.append(" | ".join(cell.strip() for cell in row))
    except csv.Error as exc:
        warnings.append(f"CSV parsing warning: {exc}")
    text = "\n".join(rows)
    status = "partial" if warnings else "success"
    return ExtractedContent(
        text=text,
        title=path.stem,
        extraction_method="csv",
        pages=[PageText(page_number=1, text=text)] if text else [],
        warnings=warnings,
        extraction_status=status,
    )


def _extract_html(path: Path, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    raw, warnings = _read_text_best_effort(path)
    if strategy == "raw":
        return ExtractedContent(
            text=raw,
            title=path.stem,
            extraction_method="html-raw",
            pages=[PageText(page_number=1, text=raw)] if raw else [],
            warnings=warnings + ["Raw HTML retry strategy used."],
            extraction_status="partial",
            fallback_chain=["html-raw"],
        )
    try:
        from bs4 import BeautifulSoup  # type: ignore
    except ImportError:
        return ExtractedContent(
            text=raw,
            title=path.stem,
            extraction_method="html-raw",
            pages=[PageText(page_number=1, text=raw)] if raw else [],
            warnings=warnings + ["BeautifulSoup not installed; HTML parsed as raw text."],
            extraction_status="partial",
            fallback_chain=["html-bs4", "html-raw"],
        )
    try:
        soup = BeautifulSoup(raw, "html.parser")
        title = soup.title.string.strip() if soup.title and soup.title.string else path.stem
        text = soup.get_text("\n", strip=True)
        return ExtractedContent(
            text=text,
            title=title,
            extraction_method="html-bs4",
            pages=[PageText(page_number=1, text=text)] if text else [],
            warnings=warnings,
        )
    except Exception as exc:
        return ExtractedContent(
            text=raw,
            title=path.stem,
            extraction_method="html-raw",
            pages=[PageText(page_number=1, text=raw)] if raw else [],
            warnings=warnings + [f"HTML parser fallback used: {exc}"],
            extraction_status="partial",
            fallback_chain=["html-bs4", "html-raw"],
            failure_reason=str(exc),
        )


def _extract_xml(path: Path, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    raw, warnings = _read_text_best_effort(path)
    if strategy == "raw":
        return ExtractedContent(
            text=raw,
            title=path.stem,
            extraction_method="xml-raw",
            pages=[PageText(page_number=1, text=raw)] if raw else [],
            warnings=warnings + ["Raw XML retry strategy used."],
            extraction_status="partial",
            fallback_chain=["xml-raw"],
        )
    try:
        root = ElementTree.fromstring(raw)
        text = "\n".join(" ".join(elem.itertext()).strip() for elem in root.iter() if " ".join(elem.itertext()).strip())
        title = root.tag
        method = "xml-tree"
        status = "success"
    except ElementTree.ParseError as exc:
        text = raw
        title = path.stem
        method = "xml-raw"
        status = "partial"
        warnings.append(f"Malformed XML parsed as raw text: {exc}")
    return ExtractedContent(
        text=text,
        title=title,
        extraction_method=method,
        pages=[PageText(page_number=1, text=text)] if text else [],
        warnings=warnings,
        extraction_status=status,
        fallback_chain=["xml-tree", method],
    )


def _extract_json(path: Path, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    raw, warnings = _read_text_best_effort(path)
    if strategy == "raw":
        return ExtractedContent(
            text=raw,
            title=path.stem,
            extraction_method="json-raw",
            pages=[PageText(page_number=1, text=raw)] if raw else [],
            warnings=warnings + ["Raw JSON retry strategy used."],
            extraction_status="partial",
            fallback_chain=["json-raw"],
        )
    try:
        obj = json.loads(raw)
        text = json.dumps(obj, indent=2, ensure_ascii=False)
        method = "json"
        status = "success"
    except json.JSONDecodeError as exc:
        text = raw
        method = "json-raw"
        status = "partial"
        warnings.append(f"Malformed JSON parsed as raw text: {exc}")
    return ExtractedContent(
        text=text,
        title=path.stem,
        extraction_method=method,
        pages=[PageText(page_number=1, text=text)] if text else [],
        warnings=warnings,
        extraction_status=status,
        fallback_chain=["json", method],
    )


def _extract_pdf(path: Path, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    warnings: list[str] = []
    fallback_chain: list[str] = []

    if strategy == "pymupdf_only":
        fitz_content = _extract_pdf_with_fitz(path, config)
        if fitz_content is None:
            return _make_metadata_only(path, "pymupdf-missing", "PyMuPDF not installed for retry strategy.")
        fitz_content.fallback_chain = ["pymupdf_only"]
        return fitz_content
    if strategy == "pdfplumber_only":
        pdfplumber_content = _extract_pdf_with_pdfplumber(path)
        if pdfplumber_content is None:
            return _make_metadata_only(path, "pdfplumber-missing", "pdfplumber not installed for retry strategy.")
        pdfplumber_content.fallback_chain = ["pdfplumber_only"]
        return pdfplumber_content
    if strategy == "pypdf_only":
        pypdf_content = _extract_pdf_with_pypdf(path)
        pypdf_content.fallback_chain = ["pypdf_only"]
        return pypdf_content

    fitz_content = _extract_pdf_with_fitz(path, config)
    fallback_chain.append("pymupdf")
    if fitz_content is not None and fitz_content.text.strip():
        fitz_content.fallback_chain = ["pymupdf"] if fitz_content.extraction_method == "pymupdf" else ["pymupdf", fitz_content.extraction_method]
        return fitz_content
    if fitz_content is not None:
        warnings.extend(fitz_content.warnings)
        if fitz_content.failure_reason:
            warnings.append(f"PyMuPDF extraction yielded no usable text: {fitz_content.failure_reason}")

    pdfplumber_content = _extract_pdf_with_pdfplumber(path)
    fallback_chain.append("pdfplumber")
    if pdfplumber_content is not None and pdfplumber_content.text.strip():
        pdfplumber_content.warnings = warnings + pdfplumber_content.warnings
        pdfplumber_content.extraction_status = "partial" if warnings else pdfplumber_content.extraction_status
        pdfplumber_content.fallback_chain = fallback_chain.copy()
        return pdfplumber_content
    if pdfplumber_content is not None:
        warnings.extend(pdfplumber_content.warnings)

    pypdf_content = _extract_pdf_with_pypdf(path)
    pypdf_content.warnings = warnings + pypdf_content.warnings
    pypdf_content.fallback_chain = fallback_chain + ["pypdf"]
    if not pypdf_content.text.strip() and pypdf_content.extraction_status == "success":
        pypdf_content.extraction_status = "metadata_only"
        pypdf_content.warnings.append("No usable PDF text was extracted; metadata-only fallback stored.")
    elif pypdf_content.text.strip() and warnings:
        pypdf_content.extraction_status = "partial"
    return pypdf_content


def _extract_pdf_with_fitz(path: Path, config: Config | None) -> ExtractedContent | None:
    try:
        import fitz  # type: ignore
    except ImportError:
        return None

    try:
        with _mupdf_messages_disabled(fitz):
            doc = fitz.open(path)
            pages: list[PageText] = []
            used_ocr = False
            method = "pymupdf"

            for index, page in enumerate(doc, start=1):
                text = page.get_text("text") or ""
                image_heavy = bool(page.get_images(full=True))
                ocr_confidence = None
                ocr_used = False
                is_scanned = False

                if _looks_like_low_quality_page(text, image_heavy):
                    is_scanned = image_heavy
                    ocr_text, ocr_confidence = _ocr_page_with_tesseract(page, config)
                    if ocr_text:
                        text = ocr_text
                        ocr_used = True
                        used_ocr = True
                        method = "pymupdf+tesseract"

                pages.append(
                    PageText(
                        page_number=index,
                        text=text,
                        ocr_used=ocr_used,
                        ocr_confidence=ocr_confidence,
                        is_scanned=is_scanned,
                    )
                )
    except Exception as exc:
        return ExtractedContent(
            text="",
            title=path.stem,
            extraction_method="pymupdf-failed",
            extraction_status="failed",
            failure_reason=str(exc),
            warnings=[f"PyMuPDF failed: {exc}"],
            quality_score=0.0,
        )

    body = "\n\n".join(page.text for page in pages if page.text.strip())
    warnings: list[str] = []
    status = "success"
    if used_ocr:
        warnings.append("OCR fallback used on one or more PDF pages.")
        status = "partial"
    if not body.strip():
        warnings.append("PyMuPDF returned no usable text.")
        status = "partial"
    return ExtractedContent(
        text=body,
        title=path.stem,
        extraction_method=method,
        pages=pages,
        ocr_used=used_ocr,
        warnings=warnings,
        extraction_status=status,
    )


@contextmanager
def _mupdf_messages_disabled(fitz_module):
    tools = getattr(fitz_module, "TOOLS", None)
    if tools is None:
        yield
        return

    previous_errors = _call_mupdf_toggle(tools.mupdf_display_errors)
    previous_warnings = _call_mupdf_toggle(tools.mupdf_display_warnings)
    try:
        tools.mupdf_display_errors(False)
        tools.mupdf_display_warnings(False)
        yield
    finally:
        tools.mupdf_display_errors(previous_errors)
        tools.mupdf_display_warnings(previous_warnings)
        reset = getattr(tools, "reset_mupdf_warnings", None)
        if callable(reset):
            reset()


def _call_mupdf_toggle(toggle) -> bool:
    try:
        return bool(toggle())
    except TypeError:
        return True


def _extract_pdf_with_pdfplumber(path: Path) -> ExtractedContent | None:
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        return None

    pages: list[PageText] = []
    try:
        with pdfplumber.open(path) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                pages.append(PageText(page_number=index, text=text))
    except Exception as exc:
        return ExtractedContent(
            text="",
            title=path.stem,
            extraction_method="pdfplumber-failed",
            extraction_status="failed",
            failure_reason=str(exc),
            warnings=[f"pdfplumber failed: {exc}"],
            quality_score=0.0,
        )
    body = "\n\n".join(page.text for page in pages if page.text.strip())
    status = "success" if body.strip() else "partial"
    warnings = [] if body.strip() else ["pdfplumber returned no usable text."]
    return ExtractedContent(
        text=body,
        title=path.stem,
        extraction_method="pdfplumber",
        pages=pages,
        warnings=warnings,
        extraction_status=status,
    )


def _extract_pdf_with_pypdf(path: Path) -> ExtractedContent:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        return _make_metadata_only(path, "pdf-missing", "Install pypdf or PyMuPDF to parse PDF.")

    try:
        reader = PdfReader(str(path))
        pages: list[PageText] = []
        for index, page in enumerate(reader.pages, start=1):
            pages.append(PageText(page_number=index, text=page.extract_text() or ""))
    except Exception as exc:
        return ExtractedContent(
            text="",
            title=path.stem,
            extraction_method="pypdf-failed",
            extraction_status="failed",
            failure_reason=str(exc),
            warnings=[f"pypdf failed: {exc}"],
            quality_score=0.0,
        )
    body = "\n\n".join(page.text for page in pages if page.text.strip())
    status = "success" if body.strip() else "metadata_only"
    warnings = [] if body.strip() else ["pypdf returned no usable text."]
    return ExtractedContent(
        text=body,
        title=path.stem,
        extraction_method="pypdf",
        pages=pages,
        warnings=warnings,
        extraction_status=status,
        quality_score=0.3 if not body.strip() else 1.0,
    )


def _extract_docx(path: Path, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    try:
        from docx import Document  # type: ignore
    except ImportError:
        return _make_metadata_only(path, "docx-missing", "Install python-docx to parse DOCX.")
    try:
        doc = Document(str(path))
        title = doc.core_properties.title or path.stem
        paragraphs = [p for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(p.text for p in paragraphs)
        structure_units: list[dict[str, str]] = []
        current_label = "Intro"
        current_lines: list[str] = []

        def flush_unit() -> None:
            if current_lines:
                structure_units.append({"label": current_label, "text": "\n".join(current_lines).strip()})

        for paragraph in paragraphs:
            style_name = str(getattr(getattr(paragraph, "style", None), "name", "") or "")
            para_text = paragraph.text.strip()
            lowered_style = style_name.lower()
            is_heading = lowered_style.startswith("heading")
            is_numbered_section = bool(re.match(r"^\d+(?:\.\d+)*\s+\S+", para_text))
            is_list = "list" in lowered_style or para_text.startswith(("- ", "* ", "1.", "2.", "3."))
            if is_numbered_section and len(para_text.split()) <= 16:
                is_heading = True
            if is_heading:
                flush_unit()
                current_label = para_text
                current_lines = []
                continue
            if is_list:
                current_lines.append(f"- {para_text.lstrip('-* ').strip()}")
            else:
                current_lines.append(para_text)
        flush_unit()

        for table_index, table in enumerate(getattr(doc, "tables", []) or [], start=1):
            rows: list[list[str]] = []
            for row in table.rows:
                values = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                if any(value for value in values):
                    rows.append(values)
            if rows:
                row_lines = _format_docx_table(rows)
                structure_units.append(
                    {
                        "label": f"Table {table_index}",
                        "text": row_lines,
                    }
                )
    except Exception as exc:
        return ExtractedContent(
            text="",
            title=path.stem,
            extraction_method="docx-failed",
            extraction_status="failed",
            failure_reason=str(exc),
            warnings=[f"DOCX extraction failed: {exc}"],
            quality_score=0.0,
        )
    status = "success" if text.strip() else "metadata_only"
    warnings = [] if text.strip() else ["DOCX contained no readable paragraph text."]
    return ExtractedContent(
        text=text,
        title=title,
        extraction_method="docx",
        pages=[PageText(page_number=1, text=text)] if text else [],
        warnings=warnings,
        extraction_status=status,
        extra={"structure_units": structure_units[:12]},
    )


def _format_docx_table(rows: list[list[str]]) -> str:
    normalized = [[cell.strip() for cell in row] for row in rows if any(cell.strip() for cell in row)]
    if not normalized:
        return ""
    column_count = max(len(row) for row in normalized)
    padded = [row + [""] * (column_count - len(row)) for row in normalized]
    widths = [max(len(row[index]) for row in padded) for index in range(column_count)]

    def format_row(row: list[str]) -> str:
        return " | ".join(cell.ljust(widths[index]) for index, cell in enumerate(row))

    header = format_row(padded[0])
    separator = "-+-".join("-" * width for width in widths)
    body = [format_row(row) for row in padded[1:]]
    return "\n".join([header, separator, *body]).strip()


def _extract_xlsx(path: Path, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError:
        return _make_metadata_only(path, "xlsx-missing", "Install openpyxl to parse XLSX.")

    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
        rows: list[str] = []
        pages: list[PageText] = []
        for index, ws in enumerate(wb.worksheets, start=1):
            sheet_rows = [f"# Sheet: {ws.title}"]
            for row in ws.iter_rows(values_only=True):
                values = ["" if v is None else str(v).strip() for v in row]
                line = " | ".join(v for v in values if v)
                if line:
                    sheet_rows.append(line)
            sheet_text = "\n".join(sheet_rows)
            rows.append(sheet_text)
            pages.append(PageText(page_number=index, text=sheet_text))
    except Exception as exc:
        return ExtractedContent(
            text="",
            title=path.stem,
            extraction_method="xlsx-failed",
            extraction_status="failed",
            failure_reason=str(exc),
            warnings=[f"XLSX extraction failed: {exc}"],
            quality_score=0.0,
        )

    text = "\n\n".join(rows)
    status = "success" if text.strip() else "metadata_only"
    warnings = [] if text.strip() else ["Workbook contained no readable values."]
    return ExtractedContent(
        text=text,
        title=path.stem,
        extraction_method="xlsx",
        pages=pages,
        warnings=warnings,
        extraction_status=status,
    )


def _extract_pptx(path: Path, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return _make_metadata_only(path, "pptx-missing", "Install python-pptx to parse PPTX.")
    try:
        prs = Presentation(str(path))
        slides: list[str] = []
        pages: list[PageText] = []
        for index, slide in enumerate(prs.slides, start=1):
            chunks = [f"# Slide {index}"]
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    chunks.append(text.strip())
            slide_text = "\n".join(chunks)
            slides.append(slide_text)
            pages.append(PageText(page_number=index, text=slide_text))
    except Exception as exc:
        return ExtractedContent(
            text="",
            title=path.stem,
            extraction_method="pptx-failed",
            extraction_status="failed",
            failure_reason=str(exc),
            warnings=[f"PPTX extraction failed: {exc}"],
            quality_score=0.0,
        )
    text = "\n\n".join(slides)
    status = "success" if text.strip() else "metadata_only"
    warnings = [] if text.strip() else ["Presentation contained no readable slide text."]
    return ExtractedContent(
        text=text,
        title=path.stem,
        extraction_method="pptx",
        pages=pages,
        warnings=warnings,
        extraction_status=status,
    )


def _extract_image(path: Path, config: Config | None = None, strategy: str | None = None) -> ExtractedContent:
    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore
    except ImportError:
        return _make_metadata_only(path, "image-ocr-missing", "Install pytesseract and Pillow to parse image documents.")

    try:
        image = Image.open(path)
        text = pytesseract.image_to_string(image).strip()
    except Exception as exc:
        return ExtractedContent(
            text="",
            title=path.stem,
            extraction_method="image-ocr-failed",
            extraction_status="failed",
            failure_reason=str(exc),
            warnings=[f"Image OCR failed: {exc}"],
            quality_score=0.0,
        )
    status = "success" if text.strip() else "partial"
    warnings = [] if text.strip() else ["OCR produced no readable text from the image."]
    return ExtractedContent(
        text=text,
        title=path.stem,
        extraction_method="image-tesseract",
        pages=[PageText(page_number=1, text=text, ocr_used=True, is_scanned=True)] if text else [PageText(page_number=1, text="", ocr_used=True, is_scanned=True)],
        ocr_used=True,
        warnings=warnings,
        extraction_status=status,
    )


def _looks_like_low_quality_page(text: str, image_heavy: bool) -> bool:
    stripped = text.strip()
    if not stripped:
        return True
    if len(stripped) < 40 and image_heavy:
        return True
    alpha = sum(1 for ch in stripped if ch.isalpha())
    return alpha < max(10, len(stripped) * 0.2) and image_heavy


def _ocr_page_with_tesseract(page, config: Config | None) -> tuple[str | None, float | None]:
    if not config or not config.ocr.enabled or config.ocr.engine.lower() != "tesseract":
        return None, None

    try:
        import fitz  # type: ignore
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore
    except ImportError:
        return None, None

    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    mode = "RGB" if pix.n >= 3 else "L"
    image = Image.frombytes(mode, [pix.width, pix.height], pix.samples)

    try:
        data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
        words = [word.strip() for word in data.get("text", []) if word.strip()]
        confidences = []
        for conf in data.get("conf", []):
            try:
                value = float(conf)
            except (TypeError, ValueError):
                continue
            if value >= 0:
                confidences.append(value)
        text = " ".join(words).strip()
        confidence = sum(confidences) / len(confidences) if confidences else None
        if not text:
            text = pytesseract.image_to_string(image).strip()
        if confidence is not None and (confidence / 100) < config.ocr.threshold:
            return None, confidence
        return text or None, confidence
    except Exception:
        return None, None
